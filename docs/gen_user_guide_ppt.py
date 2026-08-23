# -*- coding: utf-8 -*-
"""Generate cassava ERP user-guide PPT."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

OUT = Path(__file__).resolve().parent / "木薯加工厂ERP使用介绍.pptx"

C_BG = RGBColor(0xF4, 0xF7, 0xF6)
C_PRIMARY = RGBColor(0x0D, 0x6E, 0x6A)
C_ACCENT = RGBColor(0xC4, 0x5C, 0x26)
C_DARK = RGBColor(0x1A, 0x2B, 0x2A)
C_MUTED = RGBColor(0x5A, 0x6B, 0x68)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_CARD = RGBColor(0xE8, 0xF0, 0xEE)


def set_run(run, text, size=18, bold=False, color=C_DARK, font="Microsoft YaHei"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_bg(slide, prs):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_BG
    shape.line.fill.background()


def add_top_bar(slide, prs, title, subtitle=""):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_PRIMARY
    bar.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.22), Inches(12), Inches(0.7))
    tf = box.text_frame
    tf.clear()
    set_run(tf.paragraphs[0].add_run(), title, size=26, bold=True, color=C_WHITE)
    if subtitle:
        p2 = tf.add_paragraph()
        set_run(p2.add_run(), subtitle, size=13, color=RGBColor(0xC8, 0xE6, 0xE3))


def add_footer(slide, prs, page, total):
    box = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.45), Inches(12), Inches(0.35))
    set_run(
        box.text_frame.paragraphs[0].add_run(),
        f"木薯粗加工 ERP · 现场使用指南  ·  {page}/{total}",
        size=11,
        color=C_MUTED,
    )


def bullets(slide, left, top, width, height, items, size=16):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        set_run(p.add_run(), "•  " + item, size=size, color=C_DARK)


def card(slide, left, top, width, height, title, lines, accent=False):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = C_CARD if not accent else RGBColor(0xFF, 0xF0, 0xE6)
    sh.line.color.rgb = C_PRIMARY if not accent else C_ACCENT
    sh.line.width = Pt(1.25)
    tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.12), width - Inches(0.3), height - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    set_run(tf.paragraphs[0].add_run(), title, size=15, bold=True, color=C_PRIMARY if not accent else C_ACCENT)
    for line in lines:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        set_run(p2.add_run(), line, size=12, color=C_DARK)


def flow_boxes(slide, steps, y=Inches(1.5)):
    n = len(steps)
    gap = Inches(0.18)
    total_w = Inches(12.2)
    w = (total_w - gap * (n - 1)) / n
    x = Inches(0.5)
    for i, (title, sub) in enumerate(steps):
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(1.35))
        sh.fill.solid()
        sh.fill.fore_color.rgb = C_PRIMARY if i % 2 == 0 else RGBColor(0x14, 0x8A, 0x84)
        sh.line.fill.background()
        tb = slide.shapes.add_textbox(x + Inches(0.08), y + Inches(0.25), w - Inches(0.16), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), f"{i + 1}. {title}", size=13, bold=True, color=C_WHITE)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        set_run(p2.add_run(), sub, size=10, color=RGBColor(0xD5, 0xEF, 0xEC))
        if i < n - 1:
            arr = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x + w + Emu(int(gap * 0.15)),
                y + Inches(0.5),
                gap * 0.7,
                Inches(0.28),
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = C_ACCENT
            arr.line.fill.background()
        x += w + gap


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slides_meta = []

    def new_slide():
        s = prs.slides.add_slide(blank)
        add_bg(s, prs)
        slides_meta.append(s)
        return s

    # 1 cover
    s = new_slide()
    hero = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    hero.fill.solid()
    hero.fill.fore_color.rgb = C_PRIMARY
    hero.line.fill.background()
    stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.6), prs.slide_width, Inches(1.9))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = RGBColor(0x0A, 0x55, 0x52)
    stripe.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(2.5))
    tf = t.text_frame
    tf.clear()
    set_run(tf.paragraphs[0].add_run(), "木薯粗加工 ERP", size=40, bold=True, color=C_WHITE)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(12)
    set_run(p2.add_run(), "从过磅收货到产品产出 · 系统用途与使用指南", size=22, color=RGBColor(0xC8, 0xE6, 0xE3))
    p3 = tf.add_paragraph()
    p3.space_before = Pt(18)
    set_run(p3.add_run(), "看完即可上手：前置配置 · App 现场操作 · 管理端查询结算", size=14, color=RGBColor(0xA8, 0xD5, 0xD1))
    f = s.shapes.add_textbox(Inches(0.8), Inches(6.1), Inches(11), Inches(0.8))
    set_run(f.text_frame.paragraphs[0].add_run(), "适用对象：采购 / 仓管 / 产线工人 / 班组长 / 管理员", size=14, color=C_WHITE)

    # 2 toc
    s = new_slide()
    add_top_bar(s, prs, "目录", "本 PPT 带您走通整条产线")
    toc = [
        ("01", "系统解决什么问题"),
        ("02", "端到端业务主链路"),
        ("03", "关键环节细节"),
        ("04", "上线前必须完成的配置"),
        ("05", "App 使用方式（统一交互）"),
        ("06", "各角色怎么用 App"),
        ("07", "管理端做什么"),
        ("08", "常见问题与上线检查清单"),
    ]
    for i, (num, title) in enumerate(toc):
        col, row = i % 2, i // 2
        card(s, Inches(0.7) + col * Inches(6.2), Inches(1.5) + row * Inches(1.15), Inches(5.8), Inches(0.95), f"{num}  {title}", [])

    # 3 purpose
    s = new_slide()
    add_top_bar(s, prs, "系统解决什么问题", "木薯粗加工现场数字化")
    bullets(
        s,
        Inches(0.6),
        Inches(1.4),
        Inches(12),
        Inches(5.5),
        [
            "把「农户过磅 → 入厂接收 → 分箱入库 → 工序过站计件 → 成品入库」串成一条可追溯链路",
            "现场用手机 App 作业；电脑管理端负责主数据、工艺、权限、查询与结算",
            "每箱有箱码，绑定溯源码，过站扫工牌+箱码，自动累计计件产量与工钱",
            "仓前损耗、工序损耗可计量；农户结算环节可配置（入厂净重或分箱合计）",
            "目标：现场少填错、少扯皮，管理员能查清「这批货走到哪一步」",
        ],
        size=18,
    )

    # 4 main flow
    s = new_slide()
    add_top_bar(s, prs, "端到端主链路总览", "从地头到成品冷库")
    flow_boxes(
        s,
        [
            ("过磅入厂", "收货 App"),
            ("仓管接收", "仓管 App"),
            ("分箱入库", "仓管扫溯源"),
            ("工序过站", "过站 App"),
            ("成品产出", "末站入库"),
        ],
        y=Inches(1.45),
    )
    bullets(
        s,
        Inches(0.6),
        Inches(3.2),
        Inches(12),
        Inches(3.5),
        [
            "收货：扫/输溯源码建过磅单，绑定农户，推送仓管待办",
            "仓管：先确认入厂接收；再对同一溯源码分箱复磅入库（生成箱码）",
            "过站：工人扫工牌+箱码，填投料/完工重量，预览后直接过账推进工序",
            "计件：确认过站后按「完工重 × 工序工价」累计；可在「我的」核对",
            "结算：按系统设置在「入厂确认」或「分箱入库后」触发农户结算",
        ],
        size=16,
    )

    # 5 receiving
    s = new_slide()
    add_top_bar(s, prs, "环节细节① 过磅收货", "App「收货」模块")
    card(
        s,
        Inches(0.5),
        Inches(1.4),
        Inches(6.0),
        Inches(5.2),
        "现场怎么做",
        [
            "1. 打开「收货」→ 过磅入厂",
            "2. 扫或手输溯源批号",
            "3. 填毛重/扣损/净重、品种、车牌等",
            "4. 关联农户（搜索选择）",
            "5. 可拍现场照片",
            "6. 预览核对 → 确认创建并绑定",
            "7. 单据推送仓管待办",
        ],
    )
    card(
        s,
        Inches(6.8),
        Inches(1.4),
        Inches(5.9),
        Inches(5.2),
        "关键要点",
        [
            "• 溯源码与单据/农户唯一绑定",
            "• 建单即进入协同流转",
            "• 仓管可退回采购修正",
            "• 分箱入库由仓管扫溯源完成",
            "• 收货页「扫溯源分箱」可直达仓管扫码页",
        ],
        accent=True,
    )

    # 6 warehouse
    s = new_slide()
    add_top_bar(s, prs, "环节细节② 仓管接收与分箱", "App「仓管」模块")
    flow_boxes(
        s,
        [
            ("待办/扫码", "定位过磅单"),
            ("入厂接收", "核对相符确认"),
            ("再扫溯源", "分箱复磅"),
            ("箱码入库", "记仓前损耗"),
        ],
        y=Inches(1.4),
    )
    bullets(
        s,
        Inches(0.6),
        Inches(3.15),
        Inches(12),
        Inches(3.8),
        [
            "Hub 入口：扫溯源接收/分箱、待办、箱码、盘点、出入库、工序退库、库存查询",
            "入厂接收：核对票面与实物 → 预览 → 确认（本环节不分箱）",
            "分箱入库：对已入厂溯源码录入各箱复磅重量 → 预览 → 确认",
            "仓前损耗 = 票净重 − 分箱合计；系统记录扣损率；箱码自动绑定该溯源",
            "未用箱码可「销毁」；盘点/出入库支持填表预览后一次过账",
        ],
        size=15,
    )

    # 7 station
    s = new_slide()
    add_top_bar(s, prs, "环节细节③ 工序过站与计件", "App「过站」模块")
    card(
        s,
        Inches(0.5),
        Inches(1.35),
        Inches(6.0),
        Inches(5.3),
        "本人过站 / 代人过站",
        [
            "• 本人：工牌系统锁定当前用户（只读）",
            "• 代人：手输或扫描他人工牌",
            "• 箱码：手输或扫码",
            "• 填写投料重、完工重、袋数、次品类型",
            "• 下一步预览 → 确认过站直接过账",
            "• 操作人由系统按登录账号自动记录",
            "• 过站人按工牌解析；空则默认本人",
        ],
    )
    card(
        s,
        Inches(6.8),
        Inches(1.35),
        Inches(5.9),
        Inches(5.3),
        "注意",
        [
            "• 工人须在当日产线班次授权内",
            "• 卡点工序 QC 不合格会阻断过账",
            "• 班组页已去掉重复的「扫码过站」",
            "  → 统一走首页「过站」",
            "• 计件在「我的」可看今日产量工钱",
        ],
        accent=True,
    )

    # 8 workshop + mine
    s = new_slide()
    add_top_bar(s, prs, "环节细节④ 班组与「我的」", "管理例外 + 个人核对")
    card(
        s,
        Inches(0.5),
        Inches(1.4),
        Inches(6.0),
        Inches(5.2),
        "班组工作台",
        [
            "• 概览：任务/派工/今日报工",
            "• 任务、派工接收",
            "• 灵活派发 / 改派（例外场景）",
            "• 质检、废料、返修登记",
            "• 工序一览、库存查询",
            "• 剩余料退库（不回冲计件）",
            "• 正常每箱过站不在班组重复操作",
        ],
    )
    card(
        s,
        Inches(6.8),
        Inches(1.4),
        Inches(5.9),
        Inches(5.2),
        "我的",
        [
            "• 出示工牌二维码（供过站扫）",
            "• 今日产量 / 工钱核对",
            "• 消息待办",
            "• 按权限进入其它业务模块",
            "• 切换角色工作台（多角色账号）",
        ],
        accent=True,
    )

    # 9 status table
    s = new_slide()
    add_top_bar(s, prs, "关键状态一眼看懂", "协同时对得上号")
    headers = ("环节", "完成后状态", "接下来做什么")
    rows = [
        ("收货建单并推仓", "待仓管处理", "仓管在待办/扫溯源定位单据"),
        ("仓管入厂接收", "已入厂 gate_accepted", "同一溯源码可再扫做分箱"),
        ("分箱入库确认", "箱码入库 + 可结算", "按配置触发农户结算；开过站"),
        ("工序过站确认", "报工已过账 posted", "箱码进入下一工序；计件累计"),
        ("卡点 QC 不合格", "阻断过账", "不进库存过账，需返工/处理"),
    ]
    y = Inches(1.35)
    widths = (Inches(3.5), Inches(4.0), Inches(5.0))
    xs = (Inches(0.5), Inches(4.1), Inches(8.2))
    for i, h in enumerate(headers):
        sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, xs[i], y, widths[i], Inches(0.55))
        sh.fill.solid()
        sh.fill.fore_color.rgb = C_PRIMARY
        sh.line.fill.background()
        tb = s.shapes.add_textbox(xs[i] + Inches(0.1), y + Inches(0.12), widths[i] - Inches(0.2), Inches(0.4))
        set_run(tb.text_frame.paragraphs[0].add_run(), h, size=14, bold=True, color=C_WHITE)
    y += Inches(0.55)
    for ri, row in enumerate(rows):
        bg = C_CARD if ri % 2 == 0 else C_WHITE
        for i, cell in enumerate(row):
            sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, xs[i], y, widths[i], Inches(0.85))
            sh.fill.solid()
            sh.fill.fore_color.rgb = bg
            sh.line.color.rgb = RGBColor(0xD0, 0xDC, 0xD8)
            tb = s.shapes.add_textbox(xs[i] + Inches(0.1), y + Inches(0.2), widths[i] - Inches(0.2), Inches(0.55))
            tf = tb.text_frame
            tf.word_wrap = True
            set_run(tf.paragraphs[0].add_run(), cell, size=13, bold=(i == 0), color=C_DARK)
        y += Inches(0.85)

    # 10 config
    s = new_slide()
    add_top_bar(s, prs, "上线前必须完成的配置", "Admin 管理端先配好，App 才能跑")
    items = [
        ("① 组织与账号", ["部门、员工、工牌码", "用户开户、角色权限", "产线班次与当日授权"]),
        ("② 物料与仓", ["产品/品种主数据", "三仓：原料保鲜/半成品/成品", "农户档案"]),
        ("③ 工艺与工价", ["工序定义（是否计件/卡点）", "工艺流程编排并「发布」", "工序计件工价"]),
        ("④ 业务参数", ["农户结算环节 farmer_settle_point", "gate=入厂净重结算", "box_stockin=分箱后结算"]),
    ]
    for i, (title, lines) in enumerate(items):
        col, row = i % 2, i // 2
        card(
            s,
            Inches(0.5) + col * Inches(6.35),
            Inches(1.35) + row * Inches(2.7),
            Inches(6.05),
            Inches(2.45),
            title,
            ["• " + x for x in lines],
            accent=(i % 2 == 1),
        )

    # 11 config order
    s = new_slide()
    add_top_bar(s, prs, "推荐配置顺序", "按这个顺序不容易配漏")
    bullets(
        s,
        Inches(0.7),
        Inches(1.4),
        Inches(12),
        Inches(5.5),
        [
            "1. 建仓库 / 产品 / 品种 / 农户",
            "2. 建工序 → 画工艺流程 → 发布（未发布不影响现场）",
            "3. 维护工序工资（计件单价）",
            "4. 建员工并生成/绑定工牌；开 App 账号并绑角色",
            "5. 配置班次，把工人加入当日 open 班次",
            "6. 系统设置：选择农户结算环节（入厂 or 分箱后）",
            "7. 用测试溯源码跑通：收货 → 仓管接收 → 分箱 → 过站一圈",
            "8. 核对：库存箱码、过站记录、计件金额、结算单",
        ],
        size=17,
    )

    # 12 app UX
    s = new_slide()
    add_top_bar(s, prs, "App 统一使用方式", "收货 / 仓管 / 过站 同一套交互")
    flow_boxes(
        s,
        [
            ("Hub 首页", "点业务入口"),
            ("全屏子页", "盖住底栏填表"),
            ("预览核对", "有错点修改"),
            ("确认提交", "一次过账"),
        ],
        y=Inches(1.5),
    )
    bullets(
        s,
        Inches(0.6),
        Inches(3.3),
        Inches(12),
        Inches(3.5),
        [
            "码类字段（溯源码 / 箱码 / 工牌）：支持手输 + 相机扫码",
            "人员、部门等尽量下拉搜索选择，少手填 ID",
            "过站取消「提交草稿再确认」；预览无误直接提交",
            "提交成功返回 Hub，可连续作业",
        ],
        size=17,
    )

    # 13 roles
    s = new_slide()
    add_top_bar(s, prs, "各角色怎么用 App", "对号入座")
    roles = [
        ("采购/质检", ["收货：过磅入厂", "单据与任务", "看流转下一步处理人"]),
        ("仓管", ["仓管：接收与分箱", "箱码/盘点/出入库", "工序退库确认"]),
        ("计件/固定工", ["过站：本人或代人", "我的：工牌与工钱", "须在班次授权内"]),
        ("班组长", ["班组：派工/灵活派发", "质检废料返修", "不替代日常过站"]),
        ("管理员", ["电脑端配置与查询", "结算与异常处理", "权限与工艺发布"]),
    ]
    x = Inches(0.35)
    w = Inches(2.45)
    for title, lines in roles:
        card(s, x, Inches(1.4), w, Inches(5.2), title, ["• " + L for L in lines])
        x += w + Inches(0.12)

    # 14 admin
    s = new_slide()
    add_top_bar(s, prs, "电脑管理端做什么", "配置 · 查询 · 结算，不替代现场扫码")
    bullets(
        s,
        Inches(0.6),
        Inches(1.4),
        Inches(12),
        Inches(5.5),
        [
            "人事：员工、部门、工牌、开户、人事调动（下拉选人/选部门）",
            "生产：工序、工艺流程发布、过站记录查询、例外派岗",
            "库存/资产：箱码、库存、盘点单据查询",
            "采购/财务：过磅单、农户结算付款（转账单号+回单）",
            "系统设置：结算环节、审批流、基础参数",
            "原则：日常过磅/过站/分箱请用 App；管理端避免现场代录",
        ],
        size=17,
    )

    # 15 FAQ
    s = new_slide()
    add_top_bar(s, prs, "常见问题", "现场最常卡住的点")
    faqs = [
        ("扫工牌提示未授权？", "检查员工是否在当日 open 班次成员中"),
        ("找不到可过站的箱？", "是否已分箱入库；箱码是否已销毁/完工"),
        ("仓管扫溯源不能分箱？", "是否已先做入厂接收（状态须已入厂）"),
        ("计件金额为 0？", "工序是否计件、是否已发布工艺、工价是否维护"),
        ("农户何时结算？", "看系统设置 farmer_settle_point：入厂 or 分箱后"),
        ("代人过站如何记？", "操作人=登录人；过站人=扫到的工牌员工"),
    ]
    y = Inches(1.3)
    for q, a in faqs:
        tb = s.shapes.add_textbox(Inches(0.6), y, Inches(12.2), Inches(0.85))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.clear()
        set_run(tf.paragraphs[0].add_run(), "Q  " + q, size=15, bold=True, color=C_PRIMARY)
        p2 = tf.add_paragraph()
        set_run(p2.add_run(), "A  " + a, size=14, color=C_DARK)
        y += Inches(0.9)

    # 16 checklist
    s = new_slide()
    add_top_bar(s, prs, "上线检查清单", "全部勾选再开工")
    bullets(
        s,
        Inches(0.8),
        Inches(1.4),
        Inches(11.5),
        Inches(5.5),
        [
            "□ 三仓、产品、品种、农户已建",
            "□ 工艺已发布，工序工价已维护",
            "□ 员工有工牌，App 账号角色正确",
            "□ 当日班次已 open 且含产线工人",
            "□ 结算环节已按业务选定",
            "□ 测试跑通：收货→接收→分箱→过站",
            "□ 能查到箱码、过站记录、计件金额",
            "□ 仓管/采购知道退回与异常怎么处理",
        ],
        size=18,
    )

    # 17 end
    s = new_slide()
    hero = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    hero.fill.solid()
    hero.fill.fore_color.rgb = C_PRIMARY
    hero.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.5), Inches(3))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(), "记住三条", size=28, bold=True, color=C_WHITE)
    for line in [
        "1. 管理端先配置，App 再作业",
        "2. 溯源码串收货与分箱，箱码串过站",
        "3. Hub → 填表 → 预览 → 提交",
    ]:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(14)
        set_run(p2.add_run(), line, size=18, color=RGBColor(0xC8, 0xE6, 0xE3))

    total = len(slides_meta)
    for i, slide in enumerate(slides_meta):
        if i == 0 or i == total - 1:
            continue
        add_footer(slide, prs, i + 1, total)

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
